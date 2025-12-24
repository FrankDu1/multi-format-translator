"""
PPT 翻译服务 - 混合处理方案
- 直接修改可编辑文字
- OCR + Inpaint 处理图片中的文字
- 保持原始布局和格式
"""
import os
import json
import base64
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from logger_config import app_logger

# 导入配置
try:
    from config import (
        OCR_SERVICE_URL,
        INPAINT_SERVICE_URL,
        USE_INPAINT,
        UPLOAD_FOLDER
    )
except ImportError:
    app_logger.warning("⚠️ 无法导入配置，使用默认值")
    OCR_SERVICE_URL = "http://localhost:29001/ocr"
    INPAINT_SERVICE_URL = "http://localhost:29002/inpaint"
    USE_INPAINT = True
    UPLOAD_FOLDER = "uploads"

UPLOAD_DIR = os.path.join(os.path.dirname(__file__), '..', UPLOAD_FOLDER)
os.makedirs(UPLOAD_DIR, exist_ok=True)

# ============= 文本提取函数 =============

def extract_text_from_shape(shape):
    """从形状中提取文本及其格式信息"""
    text_content = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    text_content.append({
                        'text': run.text,
                        'font_size': run.font.size,
                        'font_name': run.font.name,
                        'bold': run.font.bold,
                        'italic': run.font.italic,
                        'run': run  # 保存引用以便后续更新
                    })
    return text_content


def extract_ppt_elements(ppt_path):
    """提取 PPT 的所有元素（文本、图片、图表、表格）"""
    prs = Presentation(ppt_path)
    slide_elements = []
    
    for slide_idx, slide in enumerate(prs.slides):
        app_logger.info(f"📄 处理幻灯片 {slide_idx + 1}/{len(prs.slides)}")
        
        elements = {
            'slide_idx': slide_idx,
            'texts': [],      # 可直接编辑的文本
            'images': [],     # 需要 OCR 的图片
            'charts': [],     # 图表元素
            'tables': [],     # 表格元素
            'shapes': []      # 其他形状
        }
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # 1. 处理文本框
                if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                    text_content = extract_text_from_shape(shape)
                    if text_content:
                        elements['texts'].append({
                            'shape': shape,
                            'shape_idx': shape_idx,
                            'content': text_content,
                            'type': 'text'
                        })
                
                # 2. 处理图片（可能包含文字）
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # 提取图片数据
                        image = shape.image
                        image_bytes = image.blob
                        img = Image.open(BytesIO(image_bytes))
                        
                        elements['images'].append({
                            'shape': shape,
                            'shape_idx': shape_idx,
                            'image': img,
                            'image_bytes': image_bytes,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height,
                            'type': 'image'
                        })
                        app_logger.debug(f"  图片 {shape_idx}: {img.size}")
                    except Exception as e:
                        app_logger.warning(f"  图片提取失败 (形状 {shape_idx}): {e}")
                
                # 3. 处理图表
                elif shape.has_chart:
                    chart = shape.chart
                    chart_data = {
                        'shape': shape,
                        'shape_idx': shape_idx,
                        'chart': chart,
                        'type': 'chart',
                        'title': chart.chart_title.text_frame.text if chart.has_title else '',
                        'categories': [],
                        'series': []
                    }
                    
                    # 提取图表数据
                    try:
                        if chart.plots:
                            chart_data['categories'] = [str(cat) for cat in chart.plots[0].categories]
                        for series in chart.series:
                            chart_data['series'].append({
                                'name': series.name,
                                'values': list(series.values)
                            })
                    except Exception as e:
                        app_logger.warning(f"  图表数据提取失败 (形状 {shape_idx}): {e}")
                    
                    elements['charts'].append(chart_data)
                
                # 4. 处理表格
                elif shape.has_table:
                    table = shape.table
                    table_data = {
                        'shape': shape,
                        'shape_idx': shape_idx,
                        'type': 'table',
                        'rows': []
                    }
                    
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            row_data.append({
                                'text': cell_text,
                                'cell': cell
                            })
                        table_data['rows'].append(row_data)
                    
                    elements['tables'].append(table_data)
                
                # 5. 处理分组形状
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for grouped_shape in shape.shapes:
                        if grouped_shape.has_text_frame:
                            text_content = extract_text_from_shape(grouped_shape)
                            if text_content:
                                elements['texts'].append({
                                    'shape': grouped_shape,
                                    'shape_idx': shape_idx,
                                    'content': text_content,
                                    'type': 'text_in_group'
                                })
                
            except Exception as e:
                app_logger.error(f"  处理形状失败 (幻灯片 {slide_idx}, 形状 {shape_idx}): {e}")
        
        slide_elements.append(elements)
        app_logger.info(f"  ✓ 幻灯片 {slide_idx + 1}: {len(elements['texts'])} 文本, "
                       f"{len(elements['images'])} 图片, {len(elements['charts'])} 图表, "
                       f"{len(elements['tables'])} 表格")
    
    return prs, slide_elements


# ============= 文本翻译函数 =============

def translate_text_elements(slide_elements, translator, src_lang='auto', tgt_lang='zh', enable_summary=False):
    """翻译所有文本元素"""
    app_logger.info("🔤 开始翻译文本元素...")
    # for AI summary purpose
    all_translated_text = []
    for slide_idx, elements in enumerate(slide_elements):
        app_logger.info(f"  处理幻灯片 {slide_idx + 1}...")
        
        # 1. 翻译文本框
        text_count = 0
        for text_elem in elements['texts']:
            for content in text_elem['content']:
                try:
                    original_text = content['text'].strip()
                    if original_text:
                        # 翻译
                        translated = translator.translate(original_text, src_lang, tgt_lang)
                        content['translated'] = translated

                        # for AI summary purpose
                        all_translated_text.append(translated)

                        # 🔥 关键：直接更新到原始 run 对象
                        content['run'].text = translated
                        
                        text_count += 1
                        app_logger.debug(f"    文本 {text_count}: '{original_text[:30]}' → '{translated[:30]}'")
                except Exception as e:
                    app_logger.error(f"  文本翻译失败: {e}")
                    content['translated'] = content['text']
        
        app_logger.info(f"  ✓ 翻译了 {text_count} 个文本")
        
        # 2. 翻译表格
        table_count = 0
        for table_elem in elements['tables']:
            for row_idx, row in enumerate(table_elem['rows']):
                for col_idx, cell_data in enumerate(row):
                    try:
                        original = cell_data['text'].strip()
                        if original:
                            translated = translator.translate(original, src_lang, tgt_lang)
                            cell_data['translated'] = translated
                            
                            # 🔥 直接更新单元格
                            cell_data['cell'].text = translated
                            table_count += 1
                            app_logger.debug(f"    表格[{row_idx},{col_idx}]: '{original[:20]}' → '{translated[:20]}'")
                    except Exception as e:
                        app_logger.error(f"  表格单元格翻译失败: {e}")
                        cell_data['translated'] = cell_data['text']
        
        if table_count > 0:
            app_logger.info(f"  ✓ 翻译了 {table_count} 个表格单元格")
        
        # 3. 翻译图表标题和标签
        chart_count = 0
        for chart_elem in elements['charts']:
            try:
                # 翻译标题
                if chart_elem['title']:
                    translated_title = translator.translate(chart_elem['title'], src_lang, tgt_lang)
                    chart_elem['translated_title'] = translated_title
                    
                    # 🔥 更新图表标题
                    if chart_elem['chart'].has_title:
                        chart_elem['chart'].chart_title.text_frame.text = translated_title
                        chart_count += 1
                        app_logger.debug(f"    图表标题: '{chart_elem['title'][:20]}' → '{translated_title[:20]}'")
                
                # 翻译系列名称
                for series_idx, series_data in enumerate(chart_elem['series']):
                    if series_data['name']:
                        translated_name = translator.translate(series_data['name'], src_lang, tgt_lang)
                        series_data['translated_name'] = translated_name
                        
                        # 🔥 更新系列名称
                        chart_elem['chart'].series[series_idx].name = translated_name
                        chart_count += 1
                
            except Exception as e:
                app_logger.error(f"  图表翻译失败: {e}")
        
        if chart_count > 0:
            app_logger.info(f"  ✓ 翻译了 {chart_count} 个图表元素")
        
        app_logger.info(f"  ✓ 幻灯片 {slide_idx + 1} 文本翻译完成")

    if enable_summary:
        try:
            from services.nllb_translator_pipeline import get_translator
            summary_result = None
            if all_translated_text:
                        translator = get_translator()
                        combined_text = '\n'.join(all_translated_text)
                        
                        app_logger.info(f"🧠 开始生成AI总结...")
                        
                        translation_result = translator.translate_with_summary(
                            texts=[combined_text],
                            src_lang=src_lang,
                            tgt_lang=tgt_lang,
                            enable_summary=True
                        )
                        
                        summary_result = translation_result.get('summary')
                        
                        if summary_result and summary_result.get('success'):
                            app_logger.info(f"✓ AI总结生成成功")
                        else:
                            app_logger.warning(f"⚠️ AI总结生成失败")
            return summary_result
            
        except Exception as e:
            app_logger.error(f"❌ AI总结异常: {e}")
            summary_result = {
                'success': False,
                'summary': None,
                'error': '生成总结时发生错误 🔧'
            }


# ============= OCR + Inpaint 处理函数 =============

def call_ocr_service(image_base64, src_lang='auto'):
    """调用 OCR 服务识别图片中的文字"""
    try:
        payload = {
            'image_base64': image_base64,
            'source_lang': src_lang,
            'filter_by_language': True  # 启用语言过滤
        }
        
        app_logger.debug(f"    调用 OCR: {OCR_SERVICE_URL}")
        resp = requests.post(OCR_SERVICE_URL, json=payload, timeout=60)
        
        if resp.status_code != 200:
            app_logger.error(f"    OCR 失败: HTTP {resp.status_code}")
            return None
        
        result = resp.json()
        if not result.get('success'):
            app_logger.error(f"    OCR 失败: {result.get('error')}")
            return None
        
        return result.get('result', [])
        
    except Exception as e:
        app_logger.error(f"    OCR 调用异常: {e}")
        return None


def call_inpaint_service(image_bytes, boxes, texts):
    """调用 Inpaint 服务处理图片"""
    try:
        files = {'file': ('image.jpg', image_bytes, 'image/jpeg')}
        data = {
            'boxes': json.dumps(boxes),
            'texts': json.dumps(texts)
        }
        
        app_logger.debug(f"    调用 Inpaint: {INPAINT_SERVICE_URL}")
        app_logger.debug(f"      boxes: {len(boxes)} 个区域")
        app_logger.debug(f"      texts: {len(texts)} 段文字")
        
        resp = requests.post(INPAINT_SERVICE_URL, files=files, data=data, timeout=120)
        
        if resp.status_code != 200:
            app_logger.error(f"    Inpaint 失败: HTTP {resp.status_code}")
            return None
        
        # 返回处理后的图片字节
        return resp.content
        
    except Exception as e:
        app_logger.error(f"    Inpaint 调用异常: {e}")
        return None


def process_images_with_ocr_inpaint(slide_elements, prs, translator, src_lang='auto', tgt_lang='zh'):
    """使用 OCR + Inpaint 处理图片中的文字"""
    if not USE_INPAINT:
        app_logger.info("⚠️ Inpaint 功能已禁用，跳过图片处理")
        return
    
    app_logger.info("🖼️ 开始处理图片元素...")
    
    for slide_idx, elements in enumerate(slide_elements):
        slide = prs.slides[slide_idx]
        
        for img_idx, img_elem in enumerate(elements['images']):
            try:
                app_logger.info(f"  处理图片 {img_idx + 1}/{len(elements['images'])} (幻灯片 {slide_idx + 1})")
                
                img = img_elem['image']
                image_bytes = img_elem['image_bytes']
                
                # 1. 转换为 Base64
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                
                # 2. 调用 OCR 识别
                app_logger.info(f"    OCR 识别中...")
                ocr_results = call_ocr_service(image_base64, src_lang)
                
                if not ocr_results or len(ocr_results) == 0:
                    app_logger.info(f"    未检测到文字，跳过")
                    continue
                
                # 3. 解析 OCR 结果
                boxes = []
                original_texts = []
                
                for ocr_item in ocr_results:
                    if 'res' not in ocr_item:
                        continue
                    
                    res = ocr_item['res']
                    rec_texts = res.get('rec_texts', [])
                    dt_polys = res.get('dt_polys', res.get('rec_polys', []))
                    
                    for text, poly in zip(rec_texts, dt_polys):
                        if text.strip():
                            boxes.append(poly)
                            original_texts.append(text)
                
                app_logger.info(f"    检测到 {len(original_texts)} 个文本区域")
                
                if len(boxes) == 0:
                    continue
                
                # 4. 翻译文字
                app_logger.info(f"    翻译文字...")
                translations = []
                for text in original_texts:
                    try:
                        translated = translator.translate(text, src_lang, tgt_lang)
                        translations.append(translated)
                    except Exception as e:
                        app_logger.error(f"      翻译失败: {e}")
                        translations.append(text)
                
                # 5. 构造 texts 参数（用于 Inpaint）
                texts_data = []
                for translated_text in translations:
                    texts_data.append({
                        'text': translated_text,
                        'color': [0, 0, 0],        # 黑色文字
                        'align': 'center'          # 居中对齐
                    })
                
                # 6. 调用 Inpaint 服务
                app_logger.info(f"    Inpaint 处理中...")
                processed_image_bytes = call_inpaint_service(image_bytes, boxes, texts_data)
                
                if not processed_image_bytes:
                    app_logger.warning(f"    Inpaint 处理失败，保留原图")
                    continue
                
                # 7. 替换 PPT 中的图片
                app_logger.info(f"    替换图片...")
                shape = img_elem['shape']
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                
                # 保存处理后的图片到临时文件
                temp_img_path = os.path.join(UPLOAD_DIR, f'temp_processed_{slide_idx}_{img_idx}.jpg')
                with open(temp_img_path, 'wb') as f:
                    f.write(processed_image_bytes)
                
                # 删除原图片形状
                sp = shape.element
                sp.getparent().remove(sp)
                
                # 添加新图片（保持位置和大小）
                slide.shapes.add_picture(temp_img_path, left, top, width, height)
                
                # 清理临时文件
                try:
                    os.remove(temp_img_path)
                except:
                    pass
                
                app_logger.info(f"    ✓ 图片处理完成")
                
            except Exception as e:
                app_logger.error(f"  图片处理失败 (幻灯片 {slide_idx + 1}, 图片 {img_idx + 1}): {e}")
        
        if len(elements['images']) > 0:
            app_logger.info(f"  ✓ 幻灯片 {slide_idx + 1} 图片处理完成")


# ============= 主翻译函数 =============

def translate_ppt_file(ppt_file_path, src_lang='auto', tgt_lang='zh', output_path=None, enable_summary=False):
    """主翻译函数 - PPT 混合处理"""
    try:
        app_logger.info(f"🚀 开始翻译 PPT: {ppt_file_path}")
        app_logger.info(f"   源语言: {src_lang}, 目标语言: {tgt_lang}")
        
        # 1. 提取 PPT 元素
        app_logger.info("📊 提取 PPT 内容...")
        prs, slide_elements = extract_ppt_elements(ppt_file_path)
        total_slides = len(slide_elements)
        app_logger.info(f"✅ 提取完成: {total_slides} 页")
        
        # 统计信息
        total_texts = sum(len(s['texts']) for s in slide_elements)
        total_images = sum(len(s['images']) for s in slide_elements)
        total_charts = sum(len(s['charts']) for s in slide_elements)
        total_tables = sum(len(s['tables']) for s in slide_elements)
        app_logger.info(f"📊 统计: {total_texts} 文本框, {total_images} 图片, "
                       f"{total_charts} 图表, {total_tables} 表格")
        
        # 2. 获取翻译器
        from services.nllb_translator_pipeline import get_translator
        translator = get_translator()
        summary_result = None
        # 3. 翻译文本元素（直接修改 prs 对象）
        app_logger.info("🔤 开始翻译文本...")
        if enable_summary:
            summary_result = translate_text_elements(slide_elements, translator, src_lang, tgt_lang, enable_summary)
        else:
            translate_text_elements(slide_elements, translator, src_lang, tgt_lang, enable_summary)

        # 4. 处理图片（OCR + 翻译 + Inpaint）
        if total_images > 0 and USE_INPAINT:
            app_logger.info("🖼️ 开始处理图片...")
            process_images_with_ocr_inpaint(slide_elements, prs, translator, src_lang, tgt_lang)
        else:
            app_logger.info("⏭️  跳过图片处理")
        
        # 5. 确定输出路径
        if output_path is None:
            output_path = os.path.join(UPLOAD_DIR, 'translated_' + os.path.basename(ppt_file_path))
        
        # 🔥 确保输出目录存在
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # 6. 保存翻译后的 PPT
        app_logger.info(f"💾 保存翻译后的 PPT...")
        app_logger.info(f"   输出路径: {output_path}")
        
        try:
            prs.save(output_path)
            
            # 🔥 验证文件是否真的保存成功
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
                app_logger.info(f"✅ PPT 保存成功: {output_path}")
                app_logger.info(f"   文件大小: {file_size / 1024:.1f}KB")
                
                # 🔥 再次验证：尝试重新打开
                try:
                    test_prs = Presentation(output_path)
                    app_logger.info(f"   验证通过: {len(test_prs.slides)} 页")
                except Exception as verify_error:
                    app_logger.error(f"   ⚠️ 文件验证失败: {verify_error}")
            else:
                app_logger.error(f"❌ 文件保存失败: 文件不存在")
                raise Exception(f"Failed to save PPT: {output_path}")
            
        except Exception as save_error:
            app_logger.error(f"❌ 保存 PPT 时出错: {save_error}")
            raise

        if enable_summary:
            return output_path, summary_result
        else:
            return output_path, None
        
    except Exception as e:
        app_logger.error(f"❌ PPT 翻译失败: {e}", exc_info=True)
        raise


def translate_ppt_simple(ppt_file_path, src_lang='auto', tgt_lang='zh', output_path=None, enable_summary=False):
    """简化版 PPT 翻译（仅处理文本，不处理图片）"""
    try:
        app_logger.info(f"🚀 开始简化翻译 PPT: {ppt_file_path}")
        
        from services.nllb_translator_pipeline import get_translator
        translator = get_translator()
        
        # 提取并翻译
        prs, slide_elements = extract_ppt_elements(ppt_file_path)
        summary_result = translate_text_elements(slide_elements, translator, src_lang, tgt_lang)
        
        # 确定输出路径
        if output_path is None:
            output_path = os.path.join(UPLOAD_DIR, 'translated_simple_' + os.path.basename(ppt_file_path))
        
        # 🔥 确保输出目录存在
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # 保存
        app_logger.info(f"💾 保存简化版 PPT: {output_path}")
        prs.save(output_path)
        
        # 🔥 验证
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            app_logger.info(f"✅ 简化版 PPT 保存成功")
            app_logger.info(f"   文件大小: {file_size / 1024:.1f}KB")
        else:
            raise Exception(f"Failed to save simple PPT: {output_path}")
        
        return output_path, summary_result
        
    except Exception as e:
        app_logger.error(f"❌ 简化版 PPT 翻译失败: {e}", exc_info=True)
        raise